import os
import json
import zipfile
import openai
import base64
from hashlib import sha256
import random
import string
import requests
from werkzeug.utils import secure_filename
from cryptography.fernet import Fernet
import sys

# Import necessary libraries for file processing
from pdfminer.high_level import extract_text as extract_pdf_text
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from pyxlsb import open_workbook as open_xlsb
import csv

import pandas as pd
from datetime import datetime
import logging



# Set up basic configuration for logging
logging.basicConfig(level=logging.INFO)

TOKEN_FILE_PATH = 'tokens.json'


def get_base_path():
    """
    Get the base path for resources, accounting for PyInstaller packaging.
    """
    if hasattr(sys, '_MEIPASS'):  # This attribute exists only in PyInstaller executables
        return sys._MEIPASS
    return os.path.abspath(".")

BASE_PATH = get_base_path()

prompts_folder=os.path.join(BASE_PATH, 'prompts')
encrypted_folder=os.path.join(BASE_PATH, 'encrypted')
temp_folder=os.path.join(BASE_PATH, 'temp')

################################## Encryption and access ##########################################

def generate_token(length=16):
    token = ''.join(random.choices(string.ascii_letters + string.digits, k=length))
    return token

def get_cipher(password):
    key = sha256(password.encode()).digest()
    return Fernet(base64.urlsafe_b64encode(key))

def is_upstream(org_tree, creator, current_user):
    """
    Traverse the organizational tree to check if current_user is upstream of creator.
    """
    def traverse(node, path=[]):
        for key, value in node.items():
            if key == creator or creator in value:
                return current_user in path or key == current_user
            if isinstance(value, dict):
                if traverse(value, path + [key]):
                    return True
        return False

    return traverse(org_tree["organization"])

def validate_access(filename, token, current_user):
    try:
        # Get org and user data from session ## change this in enterpeise version to use only given org_data
        org_tree_data = session.get('org_data',{current_user: {}})
        if not org_tree_data:
            raise Exception("Org tree data is missing in the session.")

        # Decrypt the file and extract metadata
        decrypted_data, metadata = decrypt_data(filename, token)
        creator = metadata["creator_signature"]

        # Check if the current user is upstream of the creator
        if is_upstream(org_tree_data, creator, current_user):
            print(f"Access granted! {current_user} is upstream of {creator}.")
            return decrypted_data
        else:
            print(f"Access denied! {current_user} is not upstream of {creator}.")
            return None
    except Exception as e:
        print(f"Access validation failed: {e}")
        return None


def encrypt_data(data, filename, tokens):

    logging.info(f"encrypt_data filename:{filename}")
    logging.info(f"encrypt_data tokens: hidden")
    # Generate a random symmetric key
    symmetric_key = Fernet.generate_key()
    cipher = Fernet(symmetric_key)
    encrypted_data = cipher.encrypt(data)
    
    # Create a dictionary to hold the encrypted symmetric key for each token
    encrypted_keys = {}
    for token in tokens:
        token_cipher = get_cipher(token)
        encrypted_key = token_cipher.encrypt(symmetric_key)
        encrypted_keys[token] = encrypted_key.decode()

    # Convert the keys dictionary to a JSON string and encode it as bytes
    keys_json = json.dumps(encrypted_keys).encode()

    # Save the encrypted data and the keys inside the file as metadata
    with open(os.path.join(encrypted_folder, filename), 'wb') as file:
        file.write(encrypted_data)
        file.write(b'\n--KEYS--\n')  # Delimiter to separate the data from the keys
        file.write(keys_json)

    logging.info(f"Horcrux saved at: {os.path.join(encrypted_folder, filename)}")
    
    return filename

def decrypt_data(filename, token):
    logging.info(filename)
    logging.info(token)
    try:
        with open(os.path.join(encrypted_folder, filename), 'rb') as file:
            content = file.read()
        
        # Split the content into the encrypted data and the keys
        encrypted_data, keys_json = content.split(b'\n--KEYS--\n')
        encrypted_keys = json.loads(keys_json.decode())
        
        # Decrypt the symmetric key using the provided token
        token_cipher = get_cipher(token)
        encrypted_symmetric_key = encrypted_keys.get(token)
        if not encrypted_symmetric_key:
            raise Exception('Token not found.')

        symmetric_key = token_cipher.decrypt(encrypted_symmetric_key.encode())
        cipher = Fernet(symmetric_key)
        decrypted_data = cipher.decrypt(encrypted_data)
        return decrypted_data
    except Exception as e:
        raise Exception('Invalid token or decryption failed.')


################################### Settings helper functions ########################################

CONFIG_PATH = os.path.expanduser('~/.myapp/config.json')

def load_config():
    """Load the config file and return it as a dictionary."""
    if os.path.exists(CONFIG_PATH):
        with open(CONFIG_PATH, 'r') as f:
            return json.load(f)
    return {}

def save_config(data):
    """Save the current session data to a JSON file."""
    os.makedirs(os.path.dirname(CONFIG_PATH), exist_ok=True)
    with open(CONFIG_PATH, 'w') as f:
        json.dump(data, f, indent=4)

################################### Token helper functions ########################################

# Helper function to save tokens to file
def save_tokens_to_file(service, tokens):
    try:
        if os.path.exists(TOKEN_FILE_PATH):
            with open(TOKEN_FILE_PATH, 'r') as file:
                all_tokens = json.load(file)
        else:
            all_tokens = {}

        all_tokens[service] = tokens

        with open(TOKEN_FILE_PATH, 'w') as file:
            json.dump(all_tokens, file)
    except Exception as e:
        print(f"Error saving tokens to file: {e}")

# Helper function to load tokens from file
def load_tokens_from_file(service):
    try:
        if os.path.exists(TOKEN_FILE_PATH):
            with open(TOKEN_FILE_PATH, 'r') as file:
                all_tokens = json.load(file)
                return all_tokens.get(service, {})
        else:
            return {}
    except Exception as e:
        print(f"Error loading tokens from file: {e}")
        return {}


################################## ONboarding ##########################################


# def load_handover(file, token):
#     temp_enc_path = os.path.join(encrypted_folder, 'temp.hrcx')
#     file.save(temp_enc_path)
#     decrypted_data = decrypt_data('temp.hrcx', token)
#     os.remove(temp_enc_path)

#     temp_zip_path = os.path.join(temp_folder, 'temp.zip')
#     with open(temp_zip_path, 'wb') as f:
#         f.write(decrypted_data)

#     with zipfile.ZipFile(temp_zip_path, 'r') as zipf:
#         zipf.extractall(temp_folder)

#     with open(os.path.join(temp_folder, 'handover_data.json'), 'r') as json_file:
#         handover_data = json.load(json_file)
    
#     os.remove(temp_zip_path)

#     # List of extracted files (excluding 'handover.json')
#     extracted_files = [f for f in os.listdir(temp_folder) if f != 'handover_data.json']
#     return handover_data, extracted_files



def load_handover(file, token):
    temp_enc_path = os.path.join(encrypted_folder, 'temp.hrcx')
    file.save(temp_enc_path)
    decrypted_data = decrypt_data('temp.hrcx', token)
    os.remove(temp_enc_path)

    temp_zip_path = os.path.join(encrypted_folder, 'temp.zip')
    with open(temp_zip_path, 'wb') as f:
        f.write(decrypted_data)

    with zipfile.ZipFile(temp_zip_path, 'r') as zipf:
        zipf.extractall(encrypted_folder)

    with open(os.path.join(encrypted_folder, 'handover_data.json'), 'r') as json_file:
        handover_data = json.load(json_file)
    
    # os.remove(temp_zip_path)

    # List of extracted files (excluding 'handover.json')
    extracted_files = [f for f in os.listdir(encrypted_folder) if f not in ['handover_data.json', 'temp.zip', 'handover.hrcx']]

    return handover_data, extracted_files


##################################### Auto Summarisation ########################################


def extract_text_from_google_doc(doc_link, google_api_key):
    # Extract the file ID and name from the link
    try:
        file_id = doc_link.split('/d/')[1].split('/')[0]
    except IndexError:
        return None, "Invalid Google Docs link."

    # Get file metadata to retrieve the file name
    metadata_url = f"https://www.googleapis.com/drive/v3/files/{file_id}?fields=name&key={google_api_key}"
    metadata_response = requests.get(metadata_url)
    if metadata_response.status_code != 200:
        return None, f"Error fetching file metadata: {metadata_response.json().get('error', {}).get('message', 'Unknown error')}"

    file_name = metadata_response.json().get('name')

    # Build the export URL
    export_url = f"https://www.googleapis.com/drive/v3/files/{file_id}/export?mimeType=text/plain&key={google_api_key}"

    # Make the request to Google Drive API
    response = requests.get(export_url)

    if response.status_code == 200:
        return {'file_name': file_name, 'link': doc_link, 'content': response.text}, None
    else:
        return None, f"Error fetching document: {response.json().get('error', {}).get('message', 'Unknown error')}"


# def extract_text_from_file(file_path, file_type):
#     """
#     Extract text from a file based on its type.
#     """
#     try:
#         if file_type == '.pdf':
#             text = extract_pdf_text(file_path)
#         elif file_type == '.docx':
#             doc = DocxDocument(file_path)
#             text = '\n'.join([para.text for para in doc.paragraphs])
#         elif file_type == '.pptx':
#             prs = Presentation(file_path)
#             text_runs = []
#             for slide in prs.slides:
#                 for shape in slide.shapes:
#                     if hasattr(shape, "text"):
#                         text_runs.append(shape.text)
#             text = '\n'.join(text_runs)
#         elif file_type == '.txt':
#             with open(file_path, 'r', encoding='utf-8') as f:
#                 text = f.read()
#         elif file_type == '.xlsx' or file_type == '.xls':
#             df = pd.read_excel(file_path)
#             buffer = []
#             buffer.append(f"Columns: {', '.join(df.columns)}")
#             buffer.append(f"Number of Rows: {len(df)}")
#             buffer.append("Basic Statistics:")
#             buffer.append(df.describe().to_string())
#             text = '\n'.join(buffer)
#         else:
#             text = None
#         return text
#     except Exception as e:
#         logging.info(f"Error extracting text from {file_path}: {e}")
#         print(f"Error extracting text from {file_path}: {e}")
#         return None

def extract_text_from_file(file_path, file_type):
    """
    Extract text from a file based on its type.
    """
    try:
        if file_type == '.pdf':
            text = extract_pdf_text(file_path)
        elif file_type == '.docx':
            doc = DocxDocument(file_path)
            text = '\n'.join([para.text for para in doc.paragraphs])
        elif file_type == '.pptx':
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            text = '\n'.join(text_runs)
        elif file_type == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        elif file_type in ['.xlsx', '.xls']:
            df = pd.read_excel(file_path)
            text = df.to_csv(index=False)
        elif file_type == '.xlsb':
            with open_xlsb(file_path) as wb:
                sheets = wb.sheets
                data_frames = [pd.read_excel(wb, sheet=sheet) for sheet in sheets]
                combined_df = pd.concat(data_frames, ignore_index=True)
                text = combined_df.to_csv(index=False)
        elif file_type == '.csv':
            with open(file_path, newline='', encoding='utf-8') as csvfile:
                reader = csv.reader(csvfile)
                text = '\n'.join([','.join(row) for row in reader])
        elif file_type == '.json':
            with open(file_path, 'r', encoding='utf-8') as jsonfile:
                data = json.load(jsonfile)
                text = json.dumps(data, indent=4)
        else:
            text = None
        return text
    except Exception as e:
        logging.error(f"Error extracting text from {file_path}: {e}")
        return None

    ############################# Summarize offline files #############################


def summarize_offline_files(files, openai_api_key):
    """
    Process uploaded files, extract text, and generate summaries.
    """
    summaries = []
    errors = []
    
    if files:
        for file in files:
            try:
                filename = file['filename']
                file_ext = os.path.splitext(filename)[1].lower()
                file_path = os.path.join(temp_folder, filename)

                logging.info(f"Processing file: {file_path}")

                # Extract text
                try:
                    text = extract_text_from_file(file_path, file_ext)
                    if len(text.split()) > 500 or len(text) > 3000:
                        text = text[:3000]  # Trim text to max 1000 characters if needed

                    if text:
                        prompt_path = os.path.join(prompts_folder, 'Handover_doc_summary_prompt.txt')
                        summary = generate_summary(prompt_path, text, openai_api_key)
                        summaries.append({
                            'file_name': filename,
                            'summary': summary
                        })
                    else:
                        raise ValueError("No text extracted or unsupported file type.")

                except Exception as e:
                    error_message = f"Failed processing {filename}: {e}"
                    logging.error(error_message)
                    errors.append(error_message)

            finally:
                # Ensure the temporary file is removed
                try:
                    # os.remove(file_path)
                    pass
                except Exception as e:
                    logging.warning(f"Error removing {file_path}: {e}")

    else:
        error_message = "No offline files uploaded."
        logging.error(error_message)
        errors.append(error_message)
        
    return summaries, errors



def generate_gpt_response(prompt, user_input, openai_api_key):
    openai.api_key = openai_api_key
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": prompt},
            {"role": "user", "content": user_input}
            ],
    )
    return response.choices[0].message['content']


def generate_summary(prompt_path, text, openai_api_key):
    # Load the base checklist
    with open(prompt_path, 'r') as file:
        base_prompt = file.read()
    
    prompt = f"{base_prompt}  :\n\n{text}"

    user_input = "None"
    response = generate_gpt_response(prompt, user_input, openai_api_key)
    summary = response.strip()
    return summary


def call_gpt_onboarding(query, data, onboarding_chat_history, openai_api_key):

    prompt_path = os.path.join(prompts_folder, 'Onboarding_chat_prompt.txt')

    with open(prompt_path, 'r') as file:
        base_prompt = file.read()
    prompt = f"{base_prompt} :\n\nRecent chat history:{json.dumps(onboarding_chat_history)} :\n\nData collected from uploaded file{json.dumps(data)}"

    user_input = query
    return generate_gpt_response(prompt, user_input, openai_api_key)
 
 

def summarize_files(session_data):
    """
    Summarizes all uploaded files and Google Docs links in the session.

    Returns:
        dict: A dictionary containing summaries and any errors encountered.
    """
    summaries = []
    errors = []

    # Retrieve uploaded files and Google Docs links from the session
    files = session_data.get('handover_data', {}).get('resources', [])
    google_docs_links = session_data.get('handover_data', {}).get('google_docs_links', [])
    google_api_key = session_data.get('google_api_key')
    llm_api_key = session_data.get('openai_api_key')

    # Summarize offline files
    if files and llm_api_key:
        try:
            file_summaries, file_errors = summarize_offline_files(files, llm_api_key)
            summaries.extend(file_summaries)
            errors.extend(file_errors)
        except Exception as e:
            error_message = f"Error processing offline files: {e}"
            logging.error(error_message)
            errors.append(error_message)

    # Summarize Google Docs links
    if google_docs_links and google_api_key and llm_api_key:
        links = google_docs_links.strip().splitlines()
        for link in links:
            try:
                doc_data, error = extract_text_from_google_doc(link, google_api_key)
                if error:
                    logging.error(f"Error processing Google Doc link {link}: {error}")
                    errors.append(error)
                    continue

                prompt_path = os.path.join(prompts_folder, 'Handover_doc_summary_prompt.txt')

                summary = generate_summary(prompt_path, doc_data['content'], llm_api_key)
                summaries.append({
                    'file_name': doc_data['file_name'],
                    'link': doc_data['link'],
                    'summary': summary
                })
            except Exception as e:
                error_message = f"Error summarizing Google Doc {link}: {e}"
                logging.error(error_message)
                errors.append(error_message)

    # Log and return the results
    logging.info(f"Summaries generated: {summaries}")
    logging.info(f"Errors encountered: {errors}")
    return {'summaries': summaries, 'errors': errors}





############################ Handover chat supporting functions #######################


### **Supporting Functions**

#### 1. `append_to_handover_chat_history`

def append_to_handover_chat_history(chat_history, role, message):
    """
    Append a message to the handover chat history.
    """
    if role == "user":
        chat_history.append({"user": message})
    elif role == "assistant":
        chat_history.append({"assistant": message})

    # chat_history.append({"role": role, "content": message})
    


def generate_dynamic_handover_chat_prompt(role, step, handover_data, handover_chat_history):
    """
    Generate a dynamic prompt for GPT based on role, step, and current handover data.

    Args:
        role (str): User role.
        step (int): Current step in the handover process.
        handover_data (dict): Current state of the handover data.

    Returns:
        str: The generated prompt.
    """
    base_checklist_path = os.path.join(prompts_folder, 'base_handover_checklist.txt')
    # base_checklist_path = "prompts/base_handover_checklist.txt"
    role_checklist_path_txt = f"{role.lower().replace(' ', '_')}_handover_checklist.txt"
    role_checklist_path = os.path.join(prompts_folder, role_checklist_path_txt)

    # Load the base checklist
    with open(base_checklist_path, 'r') as file:
        base_checklist = file.read()

    # Load the role-specific checklist
    try:
        with open(role_checklist_path, 'r') as file:
            role_checklist = file.read()
    except FileNotFoundError:
        role_checklist = "No specific checklist available for this role."

    prompt = f"""
    You are a structured assistant helping {role or "users"} with handovers. Your responsibilities are:

    1. Start with greeting, Ask the next relevant question based on current step in checklist and the existing handover_data (JSON), User is currently on step {step}.
    2. Always keep the conversation moving toward completing the checklist or resolving the user's queries.
    3. Call functions only when specific fields or actions need processing.
    4. If the user response is ambiguous, incomplete, or doesn't fit with any checklist item:
       - Ask a relevant clarifying question.
       - Guide the user toward providing a usable response.
    5. Call appropriate functions to save the User inputs and uploaded files, with context.:
       - Add/update the handover_data in a structured text, url/link information and Attach clarifications and context-specific guidance for better readability


    Checklist (with attached instructions):
    {base_checklist}

    Depending on role and if role:
    {role_checklist}

    Adhere to these rules:
    - Take it step by step, Stick to one step at a time, dont dump questions on the user
    - Consider context of what step of checklist you are on
    - Always respond concisely and contextually.
    - If the user input is irrelevant, gently redirect to the step at hand.
    - Use function calls when necessary.
    - Avoid vague responses. if uncertain, ask a clarifying question.

    Current Handover Data:
    {json.dumps(handover_data, indent=2)}

    conversation history, chronological:
    {json.dumps(handover_chat_history, indent=2)}
    """
    return prompt


def call_function_by_name(function_name, arguments, session_data):
    """
    Calls a function dynamically based on its name.

    Args:
        function_name (str): The name of the function to call.
        function_args (dict): Arguments to pass to the function.

    Returns:
        Any: The result of the called function.
    """

    # Ensure arguments are a dictionary
    if isinstance(arguments, str):
        try:
            arguments = json.loads(arguments)  # Parse the string into a dictionary
        except json.JSONDecodeError:
            raise ValueError(f"Arguments for {function_name} must be a valid JSON string or dictionary.")

    function_map = {
        "edit_or_submit_handover_field": edit_or_submit_handover_field,
        "score_handover": score_handover,
        # "summarize_files": summarize_files,
        # "finish_handover": finish_handover,
        "track_handover_progress": track_handover_progress,
        "skip_to_next_question": skip_to_next_question,
        "provide_clarification": provide_clarification,
        # "restart_handover_session": restart_handover_session,
        "validate_handover_data": validate_handover_data,
        "retrieve_handover_field": retrieve_handover_field
    }

    if function_name in function_map:
        return function_map[function_name](**arguments, session_data=session_data)
    else:
        raise ValueError(f"Function {function_name} not found.")


def score_handover(handover_data):
    """
    Score the handover data based on completeness and quality.

    Args:
        handover_data (dict): The handover data to score.

    Returns:
        dict: Scoring summary.
    """
    score = 0
    max_score = 100
    completion_criteria = [
        "basic_info.name",
        "basic_info.role",
        "projects",
        "tasks",
        "resources"
    ]

    for criterion in completion_criteria:
        value = handover_data
        for key in criterion.split('.'):
            value = value.get(key, {})
        if value:
            score += (max_score // len(completion_criteria))

    commentary = "Good work! Ensure key documents are uploaded." if score == max_score else "Some areas need more detail."
    return {"score": score, "commentary": commentary}


def finish_handover(session_data, summarize_file_option):
    """
    Create and encrypt the handover package.
    """
    logging.info(f"Finishing initiated......")
    logging.info(summarize_file_option)
        # Retrieve the uploaded files
    handover_data = session_data.get('handover_data', {})
    handover_chat = session_data.get('handover_chat_history', {})
    creator_signature = session_data.get('creator_signature', {})

    resources = session_data.get('uploaded_files',{})
    # resources = handover_data.get('uploaded_files', [])

    # # Process each file using its path
    # for resource in resources:
    #     file_path = resource.get('content')
    #     # Now, you can open each file from the file system when needed
    #     with open(file_path, 'rb') as file:
    #         file_content = file.read()
    
    links = session_data['handover_data'].get("uploaded_links", {})

    if not resources and not links and not handover_chat:
        logging.info(f"file check")
        return jsonify({'error': 'No files uploaded.'}), 400

    logging.info(f"Handover resources collected")

    # Summarize offline files
    llm_api_key = session_data.get('openai_api_key')
    summaries = []
    errors = []

    if resources and llm_api_key and summarize_file_option == 'yes':
        try:
            file_summaries, file_errors = summarize_offline_files(resources, llm_api_key)
            summaries.extend(file_summaries)
            errors.extend(file_errors)
            handover_data['summaries'] = summaries
        except Exception as e:
            error_message = f"Error processing offline files: {e}"
            logging.error(error_message)
            errors.append(error_message)


    # Generate token and add metadata
    token = generate_token()
    logging.info(f"Handover token: {token}")

    # Append the token details to the 'handover' dict
    if 'handover_token' not in handover_data:
        handover_data['handover_token'] = []
    handover_data['handover_token'] = token

    if 'timestamp' not in handover_data:
        handover_data['timestamp'] = []
    handover_data['timestamp'] = datetime.now().isoformat()

    # if 'creator_signature' not in handover_data:
    #     handover_data['creator_signature'] = []
    # # handover_data['creator_signature'] = creator_signature

    
    logging.info(f"summary collected:{summaries}")
    if 'summaries' not in handover_data:
        handover_data['summaries'] = []

    logging.info(f"Attempting handover archive..")

    # Package the handover data and resources
    encrypted_filename = create_handover_archive(handover_data, resources, creator_signature, [token])

    logging.info(f"Finished handover archive.. encrypted_filename: {encrypted_filename}")

        # Render the handover_complete.html page
    try:
        return token, encrypted_filename
    except Exception as e:
        logging.error(f"Error in finish_handover: {e}")
        # Optionally, you can re-raise the exception or handle it as needed
        raise e



def create_handover_archive(handover_data, files, creator_signature, tokens, archive_name='handover.hrcx'):
    """
    Create a ZIP archive of the handover data and encrypt it.
    """
    logging.info(f"Attempting handover_archive.. with signature {creator_signature}")

    handover_data.pop('creator_signature', None)

    handover_json = json.dumps(handover_data, indent=4)
    temp_zip_path = os.path.join(temp_folder, 'temp.zip')
    
    logging.info(f"Attempting handover_archive.. at location {temp_folder}")
    logging.info(f"Attempting handover_archive.. {handover_json}")

    try:
        # Create a temporary ZIP archive
        with zipfile.ZipFile(temp_zip_path, 'w') as zipf:
            # Add files from their paths
            for file_info in files:
                file_path = file_info['content']
                filename = file_info['filename']
                with open(file_path, 'rb') as file:
                    zipf.writestr(filename, file.read())

            # Add the handover data JSON
            zipf.writestr('handover_data.json', handover_json)
    except Exception as e:
        logging.error(f"Error creating ZIP: {e}")
        raise Exception(f"Error creating ZIP: {e}")
    
       # Encrypt and return filename
    with open(temp_zip_path, 'rb') as f:
        zipped_data = f.read()
    # os.remove(temp_zip_path)

    logging.info(f"Attempting encrypting data from {archive_name}")

    return encrypt_data(zipped_data, archive_name, tokens)



def call_gpt_handover(prompt, user_input, llm_api_key):
    """
    Call GPT API for a conversational response or function call.
    """
    openai.api_key = llm_api_key

    # Load function definitions
    function_definitions = define_gpt_functions()
    # Call the GPT API
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": prompt},
            {"role": "user", "content": user_input}
        ],
        functions=function_definitions,
        function_call="auto"  # Allows GPT to decide if a function call is needed
    )

    return response['choices'][0]['message']


################################ Additional helper functions ################################

def edit_or_submit_handover_field(field_name, field_value, session_data):
    """
    Edit or add a handover field to the session data.
    """
    user_input = "Sure, i am ready for next steps, what to do?"

    # Move to next step
    step_info = skip_to_next_question(session_data)
    current_step = step_info['current_step']

    # Determine role and generate a dynamic prompt
    role = session_data['handover_data'].get("role", "generic")
    prompt = generate_dynamic_handover_chat_prompt(role, current_step, session_data['handover_data'], session_data['handover_chat_history'])

    # Call GPT for a response
    llm_api_key = session_data.get('openai_api_key')
    gpt_response = call_gpt_handover(prompt, user_input, llm_api_key)
    
    handover_data = session_data.get('handover_data', {})
    
    # Add or update the handover field
    handover_data[field_name] = field_value
    
    # Write back to the session_data
    session_data['handover_data'] = handover_data
    
    return f"Updated {field_name} with value: {field_value}... {gpt_response.get('content', 'ready for step? ')}"


def skip_to_next_question(session_data):
    """
    Skips to the next question in the checklist based on the current step.
    """
    current_step = session_data.get('current_step', 1)
    current_step += 1
    session_data['current_step'] = current_step
    logging.info(f"Skipping to the next question. Current step: {current_step}")

    return {"current_step": current_step, "message": f"Skipped to step {current_step}."}


def track_handover_progress(step):
    total = 10  # Total steps in the checklist
    return f"{step}/{total} steps completed"


def provide_clarification(user_input, session_data):
    """
    Generates a clarifying question or statement based on ambiguous user input.
    """
    logging.info(f"User input requires clarification: {user_input}")
    return {"Could you provide more details or clarify your previous response?"}


def validate_handover_data(session_data):
    """
    Validates the current handover data against required fields.
    """
    handover_data = session_data.get('handover_data', {})
    required_fields = ["role", "title"]
    missing_fields = [field for field in required_fields if not handover_data.get(field)]

    if missing_fields:
        logging.warning(f"Validation failed. Missing fields: {missing_fields}")
        return {f"missing fields: {missing_fields}"}
    logging.info("Handover data validation passed.")
    return {f"Handover data looks good, press Finish Handover at your ease."}


def retrieve_handover_field(field_name, session_data):
    """
    Retrieves the value of a specific field in the handover data.
    """
    handover_data = session_data.get('handover_data', {})
    field_value = handover_data.get(field_name, None)
    if field_value:
        logging.info(f"Retrieved value for field '{field_name}': {field_value}")
        return f"field_value is '{field_value}'"
    else:
        logging.warning(f"Field '{field_name}' not found in handover data.")
        return f"Field '{field_name}' does not exist in the current handover data."



################################ Function Map ################################

def define_gpt_functions():
    """
    Defines the functions available for GPT to call during chat interactions.
    """
    return [
        {
            "name": "edit_or_submit_handover_field",
            "description": "Edit or submit a field in the handover data.",
            "parameters": {
                "type": "object",
                "properties": {
                    "field_name": {"type": "string", "description": "The field to edit or submit."},
                    "field_value": {"type": "string", "description": "The value to add or edit in the field."}
                },
                "required": ["field_name", "field_value"]
            }
        },
        {
            "name": "track_handover_progress",
            "description": "Tracks the user's progress in the handover process.",
            "parameters": {
                "type": "object",
                "properties": {
                    "step": {"type": "integer", "description": "The current step in the handover process."}
                },
                "required": ["step"]
            }
        },
        {
            "name": "score_handover",
            "description": "Scores the handover based on completeness, clarity, and attachments.",
            "parameters": {
                "type": "object",
                "properties": {
                    "handover_data": {"type": "object", "description": "The current handover_data under session data json containing handover data to score."}
                },
                "required": ['handover_data']
            }
        },

        {
            "name": "skip_to_next_question",
            "description": "Skip to the next question in the handover process.",
            "parameters": {
                "type": "object",
                "properties": {
                    "session_data": {"type": "object", "description": "The current session data."}
                },
                "required": ["session_data"]
            }
        },
        {
            "name": "provide_clarification",
            "description": "Ask a clarifying question to the user based on ambiguous input.",
            "parameters": {
                "type": "object",
                "properties": {
                    "user_input": {"type": "string", "description": "The ambiguous user input that requires clarification."}
                },
                "required": ["user_input"]
            }
        },
        # {
        #     "name": "restart_handover_session",
        #     "description": "Restart the handover session from the beginning.",
        #     "parameters": {}
        # },
        {
            "name": "validate_handover_data",
            "description": "Validate the current handover data for completeness.",
            "parameters": {
                "type": "object",
                "properties": {
                    "session_data": {"type": "object", "description": "The session data to validate."}
                },
                "required": ["session_data"]
            }
        },
        {
            "name": "retrieve_handover_field",
            "description": "Retrieve the value of a specific field in the handover data.",
            "parameters": {
                "type": "object",
                "properties": {
                    "field_name": {"type": "string", "description": "The name of the field to retrieve."}
                },
                "required": ["field_name"]
            }
        }
    ]








